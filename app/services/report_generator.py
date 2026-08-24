"""app/services/report_generator.py — the real implementation of PRD §6.7
(Analytics & Reporting): server-side PDF and Excel generation, styled to the
brand identity in app/branding.py rather than left at library defaults.

Kept deliberately dependency-light: reportlab for PDF (pure Python, no
headless-browser dependency to deploy), openpyxl for Excel. Both are
genuinely production-usable, not placeholders.
"""
from __future__ import annotations
import io
from collections import Counter
from datetime import datetime, timezone

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor
from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle, Paragraph,
                                Spacer, HRFlowable)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from ..branding import BRAND, severity_color, sentiment_color
from ..models import Workspace, Incident

NAVY = HexColor("#" + BRAND["navy"])
AMBER = HexColor("#" + BRAND["amber"])
SLATE = HexColor("#" + BRAND["slate"])


def _safe(text: str) -> str:
    """Escapes any text before it goes into a ReportLab Paragraph. This is
    not optional — Paragraph has its own small markup language (for bold,
    colour, etc.), so any '<', '>', or '&' in real content gets interpreted
    as broken markup instead of plain text. Real incidents come from scraped
    news/RSS/forum content and genuinely do contain raw HTML fragments (a
    literal <img alt="..." src="..." /> tag showed up in production and
    crashed PDF generation entirely before this fix existed) — this isn't a
    hypothetical edge case, it's confirmed, reproduced, real content."""
    return (str(text or "")
           .replace("&", "&amp;")
           .replace("<", "&lt;")
           .replace(">", "&gt;"))
OFFWHITE = HexColor("#" + BRAND["off_white"])


def generate_executive_summary(workspace: Workspace, incidents: list[Incident]) -> dict:
    """A genuine written executive summary — overview, insights, and
    recommendations — not just restated counts. Shared by both the PDF and
    Excel exports so the analysis is identical regardless of which format
    someone downloads. Falls back to a plain, honest templated summary if
    ANTHROPIC_API_KEY isn't configured or the call fails — a report should
    never fail to generate just because the AI layer is unavailable."""
    from collections import Counter
    from ..config import ANTHROPIC_API_KEY, CLASSIFIER_MODEL

    sev = Counter(i.severity for i in incidents)
    sent = Counter(i.sentiment for i in incidents if i.sentiment)
    platform = Counter(i.platform for i in incidents)
    tags = Counter(t for i in incidents for t in (i.tags or []))
    high_incidents = [i for i in incidents if i.severity == "HIGH"][:8]

    fallback = {
        "overview": f"{len(incidents)} mentions were tracked for {workspace.name} in this period, "
                   f"with {sev.get('HIGH', 0)} classified HIGH severity, {sev.get('MEDIUM', 0)} MEDIUM, "
                   f"and {sev.get('WATCH', 0)} WATCH.",
        "insights": f"The most active source was {platform.most_common(1)[0][0] if platform else 'none'}. "
                   f"{'Domain impersonation activity was detected — see tagged incidents below.' if tags.get('DOMAIN RISK') else 'No domain impersonation activity was detected this period.'}",
        "recommendations": "Review any HIGH severity incidents below and escalate through the Media Room "
                          "if action is required. Configure additional keywords in Settings if coverage feels incomplete.",
    }
    if not ANTHROPIC_API_KEY or not incidents:
        return fallback

    try:
        from anthropic import Anthropic
        client = Anthropic(api_key=ANTHROPIC_API_KEY)
        stats = (f"Brand: {workspace.name}\nTotal mentions: {len(incidents)}\n"
                f"Severity: {dict(sev)}\nSentiment: {dict(sent)}\n"
                f"Platforms: {dict(platform.most_common(6))}\nTags: {dict(tags)}\n"
                f"Sample HIGH severity items: " +
                "; ".join(f"[{i.platform}] {i.title[:150]}" for i in high_incidents))
        resp = client.messages.create(
            model=CLASSIFIER_MODEL, max_tokens=500,
            system="You write concise executive summaries for a brand reputation monitoring report. "
                  "Respond ONLY with a JSON object: {\"overview\": \"...\", \"insights\": \"...\", "
                  "\"recommendations\": \"...\"} — each value 2-4 plain sentences, no markdown, no bullet "
                  "characters. Be specific to the actual data given, not generic. If there is genuinely "
                  "nothing notable, say so plainly rather than inventing concern.",
            messages=[{"role": "user", "content": f"Write the executive summary for this data:\n{stats}"}],
        )
        text = "".join(b.text for b in resp.content if b.type == "text").strip()
        text = text.removeprefix("```json").removeprefix("```").removesuffix("```").strip()
        import json
        parsed = json.loads(text)
        return {"overview": str(parsed.get("overview", fallback["overview"]))[:800],
               "insights": str(parsed.get("insights", fallback["insights"]))[:800],
               "recommendations": str(parsed.get("recommendations", fallback["recommendations"]))[:800]}
    except Exception:  # noqa: BLE001 — a report must still generate even if this fails
        return fallback


# ==================================================================
# PDF REPORT
# ==================================================================
def generate_pdf_report(workspace: Workspace, incidents: list[Incident]) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, topMargin=22 * mm, bottomMargin=18 * mm,
                            leftMargin=18 * mm, rightMargin=18 * mm)
    styles = getSampleStyleSheet()

    brand_style = ParagraphStyle("Brand", parent=styles["Normal"], fontName="Helvetica-Bold",
                                 fontSize=10, textColor=AMBER, spaceAfter=2, characterSpace=2)
    title_style = ParagraphStyle("Title", parent=styles["Title"], fontName="Helvetica-Bold",
                                 fontSize=22, textColor=NAVY, spaceAfter=4)
    sub_style = ParagraphStyle("Sub", parent=styles["Normal"], fontName="Helvetica-Oblique",
                               fontSize=10, textColor=SLATE, spaceAfter=14)
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"], fontName="Helvetica-Bold",
                              fontSize=13, textColor=NAVY, spaceBefore=16, spaceAfter=8)
    body_style = ParagraphStyle("Body", parent=styles["Normal"], fontName="Helvetica",
                                fontSize=9.5, textColor=HexColor("#1F2937"), leading=13)

    story = [
        Paragraph(BRAND["name"].upper(), brand_style),
        Paragraph(f"{_safe(workspace.name)} — Monitoring Report", title_style),
        Paragraph(f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y, %H:%M UTC')} · "
                 f"{len(incidents)} mentions in range", sub_style),
        HRFlowable(width="100%", color=HexColor("#E5E7EB"), thickness=1, spaceAfter=12),
    ]

    # --- Executive summary ---
    exec_summary = generate_executive_summary(workspace, incidents)
    exec_label_style = ParagraphStyle("ExecLabel", parent=styles["Normal"], fontName="Helvetica-Bold",
                                      fontSize=9, textColor=AMBER, spaceAfter=2)
    exec_body_style = ParagraphStyle("ExecBody", parent=styles["Normal"], fontName="Helvetica",
                                     fontSize=10, textColor=HexColor("#1F2937"), leading=14.5, spaceAfter=10)
    story.append(Paragraph("Executive Summary", h2_style))
    for label, key in [("OVERVIEW", "overview"), ("KEY INSIGHTS", "insights"), ("RECOMMENDATIONS", "recommendations")]:
        story.append(Paragraph(label, exec_label_style))
        story.append(Paragraph(_safe(exec_summary[key]), exec_body_style))
    story.append(Spacer(1, 4))

    # --- Summary KPI row ---
    sev_counts = Counter(i.severity for i in incidents)
    total, high, medium, watch = len(incidents), sev_counts.get("HIGH", 0), sev_counts.get("MEDIUM", 0), sev_counts.get("WATCH", 0)
    kpi_data = [["TOTAL MENTIONS", "HIGH RISK", "MEDIUM", "WATCH"],
               [str(total), str(high), str(medium), str(watch)]]
    kpi_table = Table(kpi_data, colWidths=[42 * mm] * 4)
    kpi_table.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"), ("FONTSIZE", (0, 0), (-1, 0), 7.5),
        ("TEXTCOLOR", (0, 0), (-1, 0), SLATE), ("FONTNAME", (0, 1), (-1, 1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 1), (-1, 1), 20), ("TEXTCOLOR", (0, 1), (0, 1), NAVY),
        ("TEXTCOLOR", (1, 1), (1, 1), HexColor("#" + BRAND["severity_colors"]["HIGH"])),
        ("TEXTCOLOR", (2, 1), (2, 1), HexColor("#" + BRAND["severity_colors"]["MEDIUM"])),
        ("TEXTCOLOR", (3, 1), (3, 1), HexColor("#" + BRAND["severity_colors"]["WATCH"])),
        ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("LINEBELOW", (0, 0), (-1, 0), 0.5, HexColor("#E5E7EB")),
        ("TOPPADDING", (0, 0), (-1, 0), 4), ("BOTTOMPADDING", (0, 1), (-1, 1), 10),
    ]))
    story.append(kpi_table)

    # --- Severity & sentiment breakdown ---
    story.append(Paragraph("Severity &amp; Sentiment Breakdown", h2_style))
    sent_counts = Counter(i.sentiment for i in incidents if i.sentiment)
    breakdown_rows = [["Category", "Count", "Share"]]
    for label, count in list(sev_counts.items()) + list(sent_counts.items()):
        pct = f"{(count / total * 100):.0f}%" if total else "0%"
        breakdown_rows.append([label, str(count), pct])
    bt = Table(breakdown_rows, colWidths=[70 * mm, 30 * mm, 30 * mm])
    bt.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), NAVY), ("TEXTCOLOR", (0, 0), (-1, 0), OFFWHITE),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"), ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [OFFWHITE, HexColor("#FFFFFF")]),
        ("GRID", (0, 0), (-1, -1), 0.4, HexColor("#E5E7EB")),
        ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 8),
    ]))
    story.append(bt)

    # --- Top incidents ---
    story.append(Paragraph("Top Incidents", h2_style))
    top = sorted(incidents, key=lambda i: ({"HIGH": 0, "MEDIUM": 1, "WATCH": 2}.get(i.severity, 3), -i.reach))[:12]
    for inc in top:
        color = HexColor("#" + severity_color(inc.severity))
        title_text = _safe(inc.title[:220])
        if inc.url:
            # ReportLab's own <link> markup, applied AFTER escaping the title
            # text itself — the href attribute needs its own separate escaping
            # (quotes specifically), since it's a different context from the
            # visible text.
            safe_href = inc.url.replace('"', "%22")
            title_text = f'<link href="{safe_href}"><u color="#{BRAND["navy"]}">{title_text}</u></link>'
        row = Table([[Paragraph(f'<font color="#{severity_color(inc.severity)}"><b>{inc.severity}</b></font> '
                               f'&nbsp;&nbsp;{_safe(inc.ref)} · {_safe(inc.platform)}', body_style)],
                    [Paragraph(title_text, body_style)]], colWidths=[174 * mm])
        row.setStyle(TableStyle([
            ("LEFTPADDING", (0, 0), (-1, -1), 10), ("TOPPADDING", (0, 0), (-1, 0), 6),
            ("BOTTOMPADDING", (0, 1), (-1, 1), 8),
            ("LINEBEFORE", (0, 0), (0, -1), 2.5, color),
            ("BACKGROUND", (0, 0), (-1, -1), HexColor("#FAFAFA")),
        ]))
        story.append(row)
        story.append(Spacer(1, 4))

    story.append(Spacer(1, 16))
    story.append(HRFlowable(width="100%", color=HexColor("#E5E7EB"), thickness=1))
    story.append(Paragraph(f"{BRAND['name']} · {BRAND['tagline']}", sub_style))

    doc.build(story)
    return buf.getvalue()


# ==================================================================
# EXCEL EXPORT
# ==================================================================
def generate_excel_export(workspace: Workspace, incidents: list[Incident]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Incidents"

    headers = ["Ref", "Severity", "Sentiment", "Tags", "Title", "Platform", "Language",
              "Status", "Reach", "Source", "Posted"]
    header_fill = PatternFill(start_color=BRAND["navy"], end_color=BRAND["navy"], fill_type="solid")
    header_font = Font(color=BRAND["off_white"], bold=True, name="Calibri", size=11)
    thin = Side(style="thin", color="E5E7EB")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for col, h in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="left", vertical="center")
        cell.border = border
    ws.row_dimensions[1].height = 22
    ws.freeze_panes = "A2"

    for r, inc in enumerate(incidents, start=2):
        values = [inc.ref, inc.severity, inc.sentiment or "", ", ".join(inc.tags or []),
                 inc.title, inc.platform, inc.lang, inc.status, inc.reach, inc.source,
                 inc.posted_at.strftime("%Y-%m-%d %H:%M") if inc.posted_at else ""]
        for c, v in enumerate(values, start=1):
            cell = ws.cell(row=r, column=c, value=v)
            cell.border = border
            cell.alignment = Alignment(vertical="top", wrap_text=(c == 5))
            if c == 2:  # Severity column — brand color-coded, matching the dashboard exactly
                color = severity_color(inc.severity)
                cell.font = Font(color=color, bold=True)
            if c == 5 and inc.url:  # Title column — clickable straight through to the source
                cell.hyperlink = inc.url
                cell.font = Font(color="0563C1", underline="single")

    widths = [10, 10, 11, 22, 55, 14, 8, 11, 9, 12, 16]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # Summary sheet
    summary = wb.create_sheet("Summary")
    summary["A1"] = BRAND["name"]
    summary["A1"].font = Font(bold=True, size=16, color=BRAND["amber_dark"])
    summary["A2"] = f"{workspace.name} — Monitoring Export"
    summary["A2"].font = Font(bold=True, size=13, color=BRAND["navy"])
    summary["A3"] = f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y, %H:%M UTC')}"
    summary["A3"].font = Font(italic=True, size=10, color=BRAND["slate"])

    sev_counts = Counter(i.severity for i in incidents)
    exec_summary = generate_executive_summary(workspace, incidents)
    row = 5
    for label, key in [("Overview", "overview"), ("Key Insights", "insights"), ("Recommendations", "recommendations")]:
        c1 = summary.cell(row=row, column=1, value=label)
        c1.font = Font(bold=True, size=11, color=BRAND["amber_dark"])
        row += 1
        c2 = summary.cell(row=row, column=1, value=exec_summary[key])
        c2.font = Font(size=10.5, color="1F2937")
        c2.alignment = Alignment(wrap_text=True, vertical="top")
        summary.merge_cells(start_row=row, start_column=1, end_row=row, end_column=4)
        summary.row_dimensions[row].height = 48
        row += 2

    row += 1
    summary.cell(row=row, column=1, value="Severity").font = Font(bold=True)
    summary.cell(row=row, column=2, value="Count").font = Font(bold=True)
    for sev in ("HIGH", "MEDIUM", "WATCH"):
        row += 1
        c1 = summary.cell(row=row, column=1, value=sev)
        c1.font = Font(color=severity_color(sev), bold=True)
        summary.cell(row=row, column=2, value=sev_counts.get(sev, 0))
    summary.column_dimensions["A"].width = 60
    summary.column_dimensions["B"].width = 10

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ==================================================================
# PPTX REPORT — genuinely native PowerPoint charts (editable in
# PowerPoint, not static images), branded slide-by-slide, with the same
# executive summary and clickable links as the PDF and Excel.
# ==================================================================
PPTX_NAVY = RGBColor.from_string(BRAND["navy"])
PPTX_AMBER = RGBColor.from_string(BRAND["amber"])
PPTX_AMBER_DARK = RGBColor.from_string(BRAND["amber_dark"])
PPTX_WHITE = RGBColor.from_string("FFFFFF")
PPTX_SLATE = RGBColor.from_string(BRAND["slate"])
PPTX_INK = RGBColor.from_string("1F2937")
SEV_RGB = {"HIGH": RGBColor.from_string(BRAND["severity_colors"]["HIGH"]),
          "MEDIUM": RGBColor.from_string(BRAND["severity_colors"]["MEDIUM"]),
          "WATCH": RGBColor.from_string(BRAND["severity_colors"]["WATCH"])}

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def _pptx_bg(slide, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)  # send to back


def _pptx_text(slide, x, y, cx, cy, text, size=18, color=PPTX_INK, bold=False, italic=False,
              align=PP_ALIGN.LEFT, font="Calibri", anchor=None) -> None:
    box = slide.shapes.add_textbox(x, y, cx, cy)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.italic = italic; run.font.name = font


def _pptx_title_slide(prs, workspace: Workspace, incident_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, PPTX_NAVY)
    _pptx_text(slide, Inches(0.9), Inches(2.6), Inches(11.5), Inches(0.5),
              BRAND["name"].upper(), size=16, color=PPTX_AMBER, bold=True)
    _pptx_text(slide, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.3),
              f"{workspace.name} — Monitoring Report", size=40, color=PPTX_WHITE, bold=True)
    _pptx_text(slide, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.5),
              f"Generated {datetime.now(timezone.utc).strftime('%d %B %Y')} · {incident_count} mentions in range",
              size=15, color=RGBColor.from_string("94A3B8"), italic=True)
    return slide


def _pptx_exec_summary_slide(prs, exec_summary: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, PPTX_WHITE)
    _pptx_text(slide, Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.7),
              "Executive Summary", size=30, color=PPTX_NAVY, bold=True)
    y = Inches(1.5)
    for label, key in [("OVERVIEW", "overview"), ("KEY INSIGHTS", "insights"), ("RECOMMENDATIONS", "recommendations")]:
        _pptx_text(slide, Inches(0.7), y, Inches(11.9), Inches(0.4), label, size=13, color=PPTX_AMBER_DARK, bold=True)
        y += Inches(0.45)
        _pptx_text(slide, Inches(0.7), y, Inches(11.9), Inches(1.3), exec_summary[key], size=15, color=PPTX_INK)
        y += Inches(1.55)
    return slide


def _pptx_kpi_slide(prs, sev_counts: Counter, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, PPTX_WHITE)
    _pptx_text(slide, Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.7),
              "At a Glance", size=30, color=PPTX_NAVY, bold=True)
    cards = [("TOTAL MENTIONS", str(total), PPTX_NAVY), ("HIGH RISK", str(sev_counts.get("HIGH", 0)), SEV_RGB["HIGH"]),
            ("MEDIUM", str(sev_counts.get("MEDIUM", 0)), SEV_RGB["MEDIUM"]), ("WATCH", str(sev_counts.get("WATCH", 0)), SEV_RGB["WATCH"])]
    card_w = Inches(2.8); gap = Inches(0.3); start_x = Inches(0.7)
    for i, (label, value, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), card_w, Inches(2.2))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string("F8FAFC")
        box.line.color.rgb = RGBColor.from_string("E5E7EB"); box.line.width = Pt(1)
        box.shadow.inherit = False
        _pptx_text(slide, x + Inches(0.15), Inches(2.45), card_w - Inches(0.3), Inches(0.4),
                  label, size=11, color=PPTX_SLATE, bold=True, align=PP_ALIGN.CENTER)
        _pptx_text(slide, x + Inches(0.15), Inches(2.95), card_w - Inches(0.3), Inches(1.2),
                  value, size=48, color=color, bold=True, align=PP_ALIGN.CENTER)
    return slide


def _pptx_chart_slide(prs, title: str, categories: list, values: list, chart_type, colors: list = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_bg(slide, PPTX_WHITE)
    _pptx_text(slide, Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.7), title, size=30, color=PPTX_NAVY, bold=True)
    if not categories:
        _pptx_text(slide, Inches(0.7), Inches(3), Inches(11.5), Inches(1), "No data in this range.", size=16, color=PPTX_SLATE)
        return slide
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Mentions", values)
    gframe = slide.shapes.add_chart(chart_type, Inches(1.5), Inches(1.6), Inches(10.3), Inches(5.4), chart_data)
    chart = gframe.chart
    chart.has_legend = chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT)
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = "0"
        plot.data_labels.number_format_is_linked = False
        if colors and chart_type in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
            for i, point in enumerate(plot.series[0].points):
                if i < len(colors):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = colors[i]
        elif colors:
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = colors[0]
    except Exception:  # noqa: BLE001 — a report should still generate even if data-label styling fails
        pass
    return slide


def _pptx_top_incidents_slides(prs, incidents: list[Incident]):
    top = sorted(incidents, key=lambda i: ({"HIGH": 0, "MEDIUM": 1, "WATCH": 2}.get(i.severity, 3), -i.reach))[:12]
    per_slide = 4
    for chunk_start in range(0, len(top), per_slide):
        chunk = top[chunk_start:chunk_start + per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _pptx_bg(slide, PPTX_WHITE)
        _pptx_text(slide, Inches(0.7), Inches(0.4), Inches(11.5), Inches(0.6),
                  "Top Incidents" + (f" ({chunk_start + 1}-{chunk_start + len(chunk)})" if len(top) > per_slide else ""),
                  size=26, color=PPTX_NAVY, bold=True)
        y = Inches(1.3)
        for inc in chunk:
            color = SEV_RGB.get(inc.severity, PPTX_SLATE)
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), y, Inches(0.08), Inches(1.25))
            bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background(); bar.shadow.inherit = False
            box = slide.shapes.add_textbox(Inches(0.95), y, Inches(11.2), Inches(1.25))
            tf = box.text_frame; tf.word_wrap = True
            p1 = tf.paragraphs[0]
            r1 = p1.add_run(); r1.text = inc.severity; r1.font.bold = True; r1.font.size = Pt(13); r1.font.color.rgb = color
            r2 = p1.add_run(); r2.text = f"   {inc.ref} · {inc.platform}"; r2.font.size = Pt(12); r2.font.color.rgb = PPTX_SLATE
            p2 = tf.add_paragraph()
            r3 = p2.add_run(); r3.text = inc.title[:200]
            r3.font.size = Pt(13); r3.font.color.rgb = PPTX_INK
            if inc.url:
                r3.hyperlink.address = inc.url
                r3.font.underline = True
                r3.font.color.rgb = PPTX_NAVY
            y += Inches(1.4)
    return prs


def generate_pptx_report(workspace: Workspace, incidents: list[Incident]) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    sev_counts = Counter(i.severity for i in incidents)
    sent_counts = Counter(i.sentiment for i in incidents if i.sentiment)
    platform_counts = Counter(i.platform for i in incidents)
    exec_summary = generate_executive_summary(workspace, incidents)

    _pptx_title_slide(prs, workspace, len(incidents))
    _pptx_exec_summary_slide(prs, exec_summary)
    _pptx_kpi_slide(prs, sev_counts, len(incidents))

    sev_labels = [s for s in ("HIGH", "MEDIUM", "WATCH") if sev_counts.get(s)]
    _pptx_chart_slide(prs, "Severity Breakdown", sev_labels, [sev_counts[s] for s in sev_labels],
                      XL_CHART_TYPE.PIE, colors=[SEV_RGB[s] for s in sev_labels])

    sent_labels = list(sent_counts.keys())
    _pptx_chart_slide(prs, "Sentiment Breakdown", sent_labels, [sent_counts[s] for s in sent_labels], XL_CHART_TYPE.DOUGHNUT)

    plat_top = platform_counts.most_common(8)
    _pptx_chart_slide(prs, "Mentions by Platform", [p for p, _ in plat_top], [c for _, c in plat_top],
                      XL_CHART_TYPE.BAR_CLUSTERED, colors=[PPTX_NAVY])

    by_day: dict[str, int] = {}
    for i in incidents:
        if i.posted_at:
            key = i.posted_at.strftime("%m/%d")
            by_day[key] = by_day.get(key, 0) + 1
    trend = list(by_day.items())[-14:]
    _pptx_chart_slide(prs, "Mentions Over Time", [d for d, _ in trend], [c for _, c in trend],
                      XL_CHART_TYPE.LINE_MARKERS, colors=[PPTX_AMBER])

    _pptx_top_incidents_slides(prs, incidents)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
